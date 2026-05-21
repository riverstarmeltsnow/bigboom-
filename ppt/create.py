"""
PPT 生成工具
用法：直接运行生成示例，或 from create import PPTBuilder 在脚本中使用
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn


# ------------------------------------------------------------------
# 主题配色
# ------------------------------------------------------------------

THEMES = {
    "blue": {
        "primary":  RGBColor(0x1A, 0x27, 0x3A),   # 深蓝黑
        "accent":   RGBColor(0x29, 0x80, 0xB9),   # 蓝色
        "light":    RGBColor(0xD6, 0xEA, 0xF8),   # 浅蓝
        "warm":     RGBColor(0xE8, 0x9C, 0x31),   # 金色点缀
    },
    "green": {
        "primary":  RGBColor(0x1B, 0x43, 0x34),
        "accent":   RGBColor(0x27, 0xAE, 0x60),
        "light":    RGBColor(0xD5, 0xF5, 0xE3),
        "warm":     RGBColor(0xF3, 0x9C, 0x12),
    },
    "purple": {
        "primary":  RGBColor(0x2C, 0x1E, 0x4A),
        "accent":   RGBColor(0x8E, 0x44, 0xAD),
        "light":    RGBColor(0xE8, 0xDA, 0xEF),
        "warm":     RGBColor(0xF1, 0xC4, 0x0F),
    },
    "orange": {
        "primary":  RGBColor(0x3E, 0x27, 0x22),
        "accent":   RGBColor(0xD3, 0x54, 0x00),
        "light":    RGBColor(0xFB, 0xE9, 0xE0),
        "warm":     RGBColor(0xFE, 0xA8, 0x2F),
    },
}


def _hex_color(hex_str):
    """'#3498DB' 或 '3498DB' → RGBColor"""
    h = hex_str.strip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ------------------------------------------------------------------
# PPTBuilder
# ------------------------------------------------------------------

class PPTBuilder:
    """PPT 构建器 —— 生成带视觉风格的幻灯片"""

    def __init__(self, theme="blue"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.theme = THEMES.get(theme, THEMES["blue"])
        self._page_num = 0  # 页码计数器

    # --------------------------------------------------------------
    # 内部辅助
    # --------------------------------------------------------------

    def _add_shape(self, slide, left, top, width, height, fill_color, corner_radius=None):
        """添加矩形色块"""
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        # 设置圆角（通过 XML 属性）
        if corner_radius:
            spPr = shape._element.spPr
            prstGeom = spPr.prstGeom
            prstGeom.set("prst", "roundRect")
            avLst = prstGeom.makeelement(qn("a:avLst"), {})
            gd = avLst.makeelement(qn("a:gd"), {
                qn("a:name"): "adj",
                qn("a:fmla"): f"val {int(corner_radius)}",
            })
            avLst.append(gd)
            prstGeom.append(avLst)
        return shape

    def _add_textbox(self, slide, left, top, width, height, text,
                     font_size=18, bold=False, color=RGBColor(0x33, 0x33, 0x33),
                     align=PP_ALIGN.LEFT, font_name=None):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        if font_name:
            p.font.name = font_name
        return tf

    def _add_footer(self, slide, text="", color=RGBColor(0xBB, 0xBB, 0xBB)):
        """添加页脚（页码 + 自定义文本）"""
        self._page_num += 1
        # 底部细线
        self._add_shape(slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.02),
                        RGBColor(0xDD, 0xDD, 0xDD))
        # 页码
        self._add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1), Inches(0.4),
                          str(self._page_num), font_size=10,
                          color=color, align=PP_ALIGN.RIGHT)
        if text:
            self._add_textbox(slide, Inches(0.8), Inches(7.05), Inches(5), Inches(0.4),
                              text, font_size=10, color=color)

    # --------------------------------------------------------------
    # 幻灯片类型
    # --------------------------------------------------------------

    def add_title_slide(self, title, subtitle=None, subtitle2=None):
        """封面 —— 深色渐变背景 + 大标题"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = self._add_shape(slide, Inches(0), Inches(0),
                             Inches(13.333), Inches(7.5), self.theme["primary"])
        # 左侧装饰竖条
        self._add_shape(slide, Inches(0.8), Inches(1.5), Inches(0.08), Inches(4.5),
                        self.theme["warm"])
        # 主标题
        self._add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                          title, font_size=44, bold=True,
                          color=RGBColor(0xFF, 0xFF, 0xFF))
        # 副标题
        if subtitle:
            self._add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
                              subtitle, font_size=22,
                              color=RGBColor(0xBB, 0xCC, 0xDD))
        if subtitle2:
            self._add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.8),
                              subtitle2, font_size=16,
                              color=RGBColor(0x88, 0x99, 0xAA))
        self._page_num = 0  # 封面不计页码
        return slide

    def add_section_slide(self, title, subtitle=None):
        """章节过渡页 —— 左竖条 + 大字"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = self._add_shape(slide, Inches(0), Inches(0),
                             Inches(13.333), Inches(7.5), self.theme["primary"])
        # 中央装饰块
        self._add_shape(slide, Inches(5.5), Inches(2.8), Inches(2.3), Inches(0.06),
                        self.theme["warm"])
        self._add_shape(slide, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.06),
                        self.theme["warm"])
        # 标题
        self._add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(1.8),
                          title, font_size=40, bold=True,
                          color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        if subtitle:
            self._add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.8),
                              subtitle, font_size=18,
                              color=RGBColor(0x99, 0xAA, 0xBB), align=PP_ALIGN.CENTER)
        return slide

    def add_content_slide(self, title, items, accent_color=None):
        """内容页 —— 顶部色条 + 卡片式列表"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        ac = accent_color or self.theme["accent"]

        # 顶部装饰色条
        self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), ac)
        # 副色条
        self._add_shape(slide, Inches(0.8), Inches(0.5), Inches(0.06), Inches(0.6), ac)
        # 标题
        self._add_textbox(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.8),
                          title, font_size=30, bold=True, color=self.theme["primary"])

        # 分条绘制卡片式内容
        y_start = 1.6
        card_h = 0.9
        gap = 0.15
        for i, item in enumerate(items):
            y = y_start + i * (card_h + gap)
            # 卡片背景
            color_idx = i % 4
            card_colors = [
                RGBColor(0xF0, 0xF4, 0xF8),
                RGBColor(0xF8, 0xF0, 0xF0),
                RGBColor(0xF0, 0xF8, 0xF4),
                RGBColor(0xF8, 0xF4, 0xF0),
            ]
            self._add_shape(slide, Inches(1.0), Inches(y), Inches(11.3), card_h,
                            card_colors[color_idx], corner_radius=12700)
            # 左侧色标
            self._add_shape(slide, Inches(1.0), Inches(y), Inches(0.06), card_h, ac)
            # 序号圆
            self._add_shape(slide, Inches(1.4), Inches(y + 0.25), Inches(0.4), Inches(0.4),
                            ac)
            num_text = self._add_textbox(slide, Inches(1.4), Inches(y + 0.22),
                                         Inches(0.4), Inches(0.4),
                                         str(i + 1), font_size=14, bold=True,
                                         color=RGBColor(0xFF, 0xFF, 0xFF),
                                         align=PP_ALIGN.CENTER)
            # 文本
            self._add_textbox(slide, Inches(2.1), Inches(y + 0.15),
                              Inches(10), Inches(0.6),
                              item, font_size=18, color=RGBColor(0x33, 0x33, 0x33))

        self._add_footer(slide, color=RGBColor(0xBB, 0xBB, 0xBB))
        return slide

    def add_two_column_slide(self, title, left_items, right_items, accent_color=None):
        """双栏内容页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        ac = accent_color or self.theme["accent"]

        # 顶部装饰色条
        self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), ac)
        self._add_shape(slide, Inches(0.8), Inches(0.5), Inches(0.06), Inches(0.6), ac)
        self._add_textbox(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.8),
                          title, font_size=30, bold=True, color=self.theme["primary"])

        def _draw_column(col_items, x_offset, col_title=None):
            """画一列卡片"""
            if col_title:
                self._add_textbox(slide, Inches(x_offset), Inches(1.3),
                                  Inches(5.5), Inches(0.5),
                                  col_title, font_size=20, bold=True,
                                  color=ac)
            y_start = 1.9
            card_h = 0.85
            gap = 0.15
            for i, item in enumerate(col_items):
                y = y_start + i * (card_h + gap)
                c_bg = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else RGBColor(0xFA, 0xFA, 0xFA)
                self._add_shape(slide, Inches(x_offset), Inches(y),
                                Inches(5.5), card_h, c_bg, corner_radius=12700)
                # 左侧小色条
                self._add_shape(slide, Inches(x_offset), Inches(y),
                                Inches(0.05), card_h, ac)
                self._add_textbox(slide, Inches(x_offset + 0.4), Inches(y + 0.15),
                                  Inches(5), Inches(0.55),
                                  item, font_size=17, color=RGBColor(0x33, 0x33, 0x33))
                # 如果是第一行，加个"01" "02"编号
                num_text = f"0{i+1}"
                self._add_textbox(slide, Inches(x_offset + 4.5), Inches(y + 0.15),
                                  Inches(0.8), Inches(0.55),
                                  num_text, font_size=11, bold=True,
                                  color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.RIGHT)

        _draw_column(left_items, 0.8)
        # 中间分隔竖线
        self._add_shape(slide, Inches(6.5), Inches(1.5), Inches(0.02), Inches(5.0),
                        RGBColor(0xE0, 0xE0, 0xE0))
        _draw_column(right_items, 6.9)

        self._add_footer(slide, color=RGBColor(0xBB, 0xBB, 0xBB))
        return slide

    def add_card_slide(self, title, cards, accent_color=None):
        """
        卡片式布局 —— 每张卡片包含标题和描述
        cards: [("标题1", "描述1"), ("标题2", "描述2"), ...]
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        ac = accent_color or self.theme["accent"]

        # 顶部装饰
        self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), ac)
        self._add_shape(slide, Inches(0.8), Inches(0.5), Inches(0.06), Inches(0.6), ac)
        self._add_textbox(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.8),
                          title, font_size=30, bold=True, color=self.theme["primary"])

        cols = 3
        card_w = 3.5
        card_h = 2.5
        gap_x = 0.5
        gap_y = 0.5
        start_x = 0.8
        start_y = 1.5

        for i, (card_title, card_desc) in enumerate(cards):
            col = i % cols
            row = i // cols
            x = start_x + col * (card_w + gap_x)
            y = start_y + row * (card_h + gap_y)

            # 卡片背景
            self._add_shape(slide, Inches(x), Inches(y),
                            Inches(card_w), card_h,
                            RGBColor(0xF8, 0xF9, 0xFA), corner_radius=12700)
            # 卡片顶部色条
            self._add_shape(slide, Inches(x), Inches(y),
                            Inches(card_w), Inches(0.06), ac)
            # 卡片标题
            self._add_textbox(slide, Inches(x + 0.3), Inches(y + 0.3),
                              Inches(card_w - 0.6), Inches(0.5),
                              card_title, font_size=22, bold=True, color=ac)
            # 卡片描述
            self._add_textbox(slide, Inches(x + 0.3), Inches(y + 1.0),
                              Inches(card_w - 0.6), Inches(1.3),
                              card_desc, font_size=15,
                              color=RGBColor(0x66, 0x66, 0x66))

        self._add_footer(slide, color=RGBColor(0xBB, 0xBB, 0xBB))
        return slide

    def add_image_slide(self, title, image_path):
        """
        图文页 —— 左侧图片，右侧标题+描述
        暂时留空，导入图片后再实现
        """
        # 未来扩展
        pass

    # --------------------------------------------------------------
    # 保存
    # --------------------------------------------------------------

    def save(self, filename):
        # 默认保存到 ppt/ 目录（脚本所在目录），不是运行目录
        import os
        base = os.path.dirname(os.path.abspath(__file__))
        path = os.path.join(base, filename) if not os.path.isabs(filename) else filename
        self.prs.save(path)
        print(f"[OK] PPT 已保存: {path}")


# 使用示例（取消注释即可运行）：
# if __name__ == "__main__":
#     ppt = PPTBuilder(theme="blue")
#     ppt.add_title_slide("标题", "副标题")
#     ppt.add_content_slide("内容页", ["第一点", "第二点"])
#     ppt.save("output.pptx")
